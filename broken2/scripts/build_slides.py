"""Build a PowerPoint deck for the Adaptive Robust CBF project.

Sources:
  - probabilistic_robotics_final_project.pdf (paper-quality writeup)
  - Maps to ~15 slides for a class presentation.

Output: docs/RL_CBF_presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path("/Users/chrisliang8/Desktop/safety-go2/docs/RL_CBF_presentation.pptx")

# Color palette (clean academic look)
NAVY = RGBColor(0x1F, 0x2A, 0x44)
TEAL = RGBColor(0x0F, 0x76, 0x84)
ACCENT = RGBColor(0xC4, 0x57, 0x2C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x10, 0x10, 0x10)


def add_title_slide(prs, title, subtitle, author):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Navy band on top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.0))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.4), prs.slide_width - Inches(1.2), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    r2.font.color.rgb = TEAL
    r2.font.italic = True

    p3 = tf.add_paragraph()
    p3.space_before = Pt(20)
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = author
    r3.font.size = Pt(16)
    r3.font.color.rgb = GRAY


def add_section_title(slide, title, color=NAVY):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), prs.slide_width - Inches(1.0), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = color

    # Underline accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(2.0), Emu(38000)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL


def add_bullet_slide(prs, title, bullets, sub=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, title)

    top = Inches(1.2)
    if sub:
        sub_tb = slide.shapes.add_textbox(Inches(0.5), top, prs.slide_width - Inches(1.0), Inches(0.6))
        stf = sub_tb.text_frame
        stf.word_wrap = True
        p = stf.paragraphs[0]
        r = p.add_run()
        r.text = sub
        r.font.size = Pt(16)
        r.font.italic = True
        r.font.color.rgb = GRAY
        top = Inches(1.85)

    body = slide.shapes.add_textbox(
        Inches(0.5), top, prs.slide_width - Inches(1.0), prs.slide_height - top - Inches(0.4)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            lead, rest = item
        else:
            lead, rest = None, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        if lead is not None:
            r1 = p.add_run()
            r1.text = lead + " "
            r1.font.size = Pt(20)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(20)
            r2.font.color.rgb = BLACK
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(20)
            r.font.color.rgb = BLACK
    return slide


def add_table(slide, header, rows, left, top, width, height,
              header_fill=NAVY, header_color=WHITE, font_size=14,
              first_col_bold=False, accent_rows=None):
    n_rows = len(rows) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    accent_rows = accent_rows or {}
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(font_size)
        r.font.bold = True
        r.font.color.rgb = header_color
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            if i in accent_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent_rows[i]
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            if first_col_bold and j == 0:
                r.font.bold = True
            r.font.color.rgb = BLACK
    return table


def add_text_block(slide, text, left, top, width, height, size=18, color=BLACK,
                   bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_callout_box(slide, text, left, top, width, height, fill_color=TEAL,
                    text_color=WHITE, size=16, bold=True):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color


# =====================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---------- Slide 1: Title ----------
add_title_slide(
    prs,
    "Adaptive Robustness Margins for CBFs",
    "via Privileged-Observation Reinforcement Learning — A Quadrupedal Safety-Filter Case Study",
    "Xianmai Liang",
)

# ---------- Slide 2: Motivation ----------
add_bullet_slide(
    prs,
    "The problem with safe robots",
    [
        ("CBFs in theory.", "Reduce safety to a pointwise affine constraint on the control input — solve a small QP, get formal forward-invariance of the safe set."),
        ("Reality breaks every assumption.", "Friction varies, payloads shift the COM, perception is noisy, the robot has a body."),
        ("Robust CBF variants add tightening parameters", "(α, φ, a, c) — but worst-case-derived values are crushingly conservative. The robot freezes or refuses to move."),
        ("Hand-tuning is the practical workaround.", "But it's brittle to environments unlike the tuning set, and gives no recipe for new scenes."),
    ],
)

# ---------- Slide 3: Research question ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Research question")
add_text_block(
    slide,
    "Can a learned policy adapt the robust-CBF parameters online,",
    Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.8),
    size=28, bold=True, color=NAVY,
)
add_text_block(
    slide,
    "conditioning on what the environment currently looks like?",
    Inches(0.6), Inches(2.3), Inches(12.0), Inches(0.8),
    size=28, bold=True, color=NAVY,
)
add_callout_box(
    slide,
    "Approach: borrow the Rapid Motor Adaptation (RMA) architecture — a privileged-observation teacher reads dynamics + a top-down occupancy grid, and outputs the robustification tuple (α, φ, a, c) per timestep, feeding a robust-CBF QP that wraps an off-the-shelf locomotion controller.",
    Inches(0.6), Inches(3.6), Inches(12.0), Inches(2.0),
    fill_color=TEAL, size=18, bold=False,
)
add_text_block(
    slide,
    "Non-negativity is enforced by construction → qualitative robust-safety guarantee preserved for whatever uncertainty bounds the policy chooses.",
    Inches(0.6), Inches(5.9), Inches(12.0), Inches(1.0),
    size=16, italic=True, color=GRAY,
)

# ---------- Slide 4: CBF formulation ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Background: control barrier functions")
add_text_block(
    slide,
    "Safe set  C = { x : h(x) ≥ 0 }.  A CBF condition tightens to a pointwise affine constraint on the control:",
    Inches(0.6), Inches(1.3), Inches(12.0), Inches(0.7),
    size=18,
)
add_callout_box(
    slide,
    "L_f h(x)  +  L_g h(x) · u   ≥   −α · h(x)",
    Inches(2.5), Inches(2.3), Inches(8.3), Inches(0.9),
    fill_color=NAVY, size=24, bold=True,
)
add_text_block(
    slide,
    "Robust variants (Kolathaya 2018 ISSf, Cosner 2022 TISSf, Dean 2019, Molnar 2021) tighten the RHS:",
    Inches(0.6), Inches(3.4), Inches(12.0), Inches(0.6),
    size=18,
)
add_callout_box(
    slide,
    "L_g h · û  −  φ ‖L_g h‖²  −  a  −  b ‖û‖₂  +  α (h − c)   ≥   0",
    Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.9),
    fill_color=ACCENT, size=20, bold=True,
)
add_text_block(
    slide,
    "Each tightening parameter absorbs a specific class of uncertainty.  Set them too large → frozen robot.  Set them too small → safety guarantee shrinks.",
    Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.5),
    size=18, italic=True, color=GRAY,
)

# ---------- Slide 5: Four-parameter mapping ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "The four-parameter action space")
add_text_block(
    slide,
    "Each parameter maps to a specific uncertainty class — the policy adapts them per timestep.",
    Inches(0.6), Inches(1.2), Inches(12.0), Inches(0.5),
    size=16, italic=True, color=GRAY,
)
add_table(
    slide,
    ["Param", "Range", "Role", "Uncertainty class"],
    [
        ["α", "[0.1, 5.0]", "Class-K slope on shifted h(x) − c", "Model / tracking error"],
        ["φ", "[0.0, 5.0]", "Coefficient on ‖L_g h‖²", "Actuation uncertainty (ISSf)"],
        ["a", "[0.0, 0.5]", "Additive RHS slack", "Measurement uncertainty"],
        ["c", "[0.0, 0.5]", "Inward shift of safe-set boundary", "Boundary correction"],
        ["b", "(unused)", "Reserved for input-dep slack b‖u‖", "Would require SOCP"],
    ],
    Inches(0.6), Inches(1.95), Inches(12.0), Inches(3.5),
    font_size=16, first_col_bold=True,
)
add_callout_box(
    slide,
    "Output non-negativity → robust-safety theorem (Cosner 2022) holds for whatever uncertainty bounds the policy's outputs encode.",
    Inches(0.6), Inches(5.7), Inches(12.0), Inches(1.2),
    fill_color=GREEN, size=16, bold=False,
)

# ---------- Slide 6: Method overview ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Method: privileged-obs teacher → robust-CBF QP")

# Pipeline boxes
y = Inches(1.6)
box_h = Inches(0.9)
box_w = Inches(2.4)

boxes = [
    ("Privileged obs\n8207-D", NAVY),
    ("CNN+MLP encoder\n→ Z(12)", TEAL),
    ("π_teacher\n→ (α, φ, a, c)", ACCENT),
    ("Robust-CBF QP\n(closed-form)", GREEN),
    ("Locomotion\n(frozen)", GRAY),
]
gap = Inches(0.05)
total_w = Inches(13.333) - Inches(1.0)
single_w = (total_w - 4 * gap) / 5

for i, (txt, color) in enumerate(boxes):
    x = Inches(0.5) + i * (single_w + gap)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, single_w, box_h)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if i < len(boxes) - 1:
        # arrow
        ax = x + single_w
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(0.25), gap, Inches(0.4))
        arrow.line.fill.background()
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLACK

# Below pipeline: bullets
add_text_block(
    slide,
    "Drop-in safety filter:  the locomotion controller is fixed and frozen — works with any user's stack.",
    Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.5),
    size=16, italic=True, color=GRAY,
)

bullets = [
    "Single-integrator surrogate: ẋ_robot = u   →   QP is a closed-form half-space projection on GPU (no infeasibilities observed).",
    "Robot footprint enters as Minkowski expansion by 0.15 m so h = 0 aligns with physical contact.",
    "u_des never enters the network — it is a sidecar to the CBF, only used in reward (‖u_safe − u_des‖²).",
    "Per-step adaptation distinguishes us from per-episode α-calibration baselines (e.g., SHIELD).",
]
body = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.3), Inches(3.7))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run()
    r.text = "• " + b
    r.font.size = Pt(17)
    r.font.color.rgb = BLACK

# ---------- Slide 7: Privileged observation + network ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Privileged observation + network architecture")

# Two-column layout
add_text_block(
    slide,
    "Privileged observation (8207-D)",
    Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.5),
    size=20, bold=True, color=NAVY,
)
add_table(
    slide,
    ["Block", "Dim", "Content"],
    [
        ["Dynamics", "15", "Body velocities, ext. force/torque, friction, COM offset, last action"],
        ["Occupancy grid", "2 × 64 × 64", "Two-frame top-down, 0.1 m / cell, 6.4 m FOV"],
    ],
    Inches(0.5), Inches(1.85), Inches(6.0), Inches(2.0),
    font_size=12, first_col_bold=True,
)
add_text_block(
    slide,
    "Two-frame stacking encodes obstacle motion (grid_t − grid_{t−1}).",
    Inches(0.5), Inches(4.0), Inches(6.0), Inches(0.7),
    size=14, italic=True, color=GRAY,
)
add_text_block(
    slide,
    "Network (~600 k actor params)",
    Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
    size=20, bold=True, color=NAVY,
)
arch = [
    "Grid:  Conv(2→16, s=2) → Conv(16→32, s=2) → Linear(8192→64)",
    "Dynamics:  Linear(15 → 64)",
    "Concat → Linear(128 → 128) → Linear(128 → z)",
    "Latent:  Z ∈ ℝ¹²  (RMA-style bottleneck; get_z(obs) exposed)",
    "π-head:  Z → (α, φ, a, b, c), tanh-squash to physical range",
    "V-head: separate critic encoder (no shared params)",
]
body = slide.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(6.0), Inches(5.0))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(arch):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "• " + b
    r.font.size = Pt(14)
    r.font.color.rgb = BLACK

# ---------- Slide 8: Reward + training ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Reward design + training")
add_text_block(
    slide,
    "Reward stack (5 terms)",
    Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.4),
    size=18, bold=True, color=NAVY,
)
add_table(
    slide,
    ["Term", "Weight", "Trigger"],
    [
        ["collision",       "−100",                          "obstacle contact (terminal)"],
        ["infeasibility",   "−10 / step",                    "QP infeasible"],
        ["u_safe deviation","−0.1 ‖u_safe − u_des‖²",        "per-step (light-touch filter)"],
        ["proximity",       "−1.0 · exp(−min_sdf / 0.5)",    "per-step (gradient near obstacles)"],
        ["action_rate",     "−5×10⁻³ ‖Δa‖²",                 "per-step (smooth CBF params)"],
    ],
    Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.4),
    font_size=12, first_col_bold=True,
)
add_text_block(
    slide,
    "Negative results we recorded",
    Inches(0.5), Inches(5.2), Inches(6.0), Inches(0.4),
    size=14, bold=True, color=ACCENT,
)
add_text_block(
    slide,
    "•  Proximity = −5 → policy maxed distance instead of progress (φ̄ rose to 2.79). Halved to −1.\n•  Weight decay 1e-4 → action_std collapsed to 0.06; exploration died. Final: 1e-5 + entropy 5e-3.",
    Inches(0.5), Inches(5.55), Inches(6.0), Inches(1.6),
    size=12, italic=True, color=GRAY,
)

add_text_block(
    slide,
    "Training (PPO)",
    Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
    size=18, bold=True, color=NAVY,
)
training = [
    "4096 parallel envs × 3000 iters  (~3 h on RTX 5090)",
    "AdamW, weight_decay = 1e-5",
    "entropy_coef = 5e-3  (load-bearing)",
    "Episode length 20 s",
    "",
    "Domain randomization (per-reset)",
    "  •  friction: μs ∈ [0.30, 1.20], μd ∈ [0.20, 1.00]",
    "  •  external force ±10 N, torque ±2 Nm",
    "  •  COM offset ±5 cm xy / ±3 cm z",
    "",
    "Multi-planner mix (the deploy-realistic distribution)",
    "  •  smooth_goal 0.40 / waypoint 0.25 / mpc 0.20",
    "  •  walk 0.05 / adversarial 0.05 / legacy 0.05",
]
body = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(6.0), Inches(5.5))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(training):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(2)
    r = p.add_run()
    r.text = b
    r.font.size = Pt(13)
    r.font.color.rgb = BLACK if not b.startswith("Domain") and not b.startswith("Multi") else NAVY
    if b.startswith("Domain") or b.startswith("Multi"):
        r.font.bold = True

# ---------- Slide 9: Evaluation setup ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Evaluation: 24-config baseline sweep")
add_text_block(
    slide,
    "Each baseline is a special case of the same robust-CBF constraint.  We tune each family's grid and report the BEST configuration per task.",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6),
    size=16, italic=True, color=GRAY,
)
add_table(
    slide,
    ["Family", "Form", "Grid", "# configs"],
    [
        ["B0  (plain)",  "φ = a = c = 0;  α only",                          "α ∈ {0.5, 1.5, 3.0}",                                  "3"],
        ["B1  (ISSf)",   "φ scalar, a = c = 0",                              "α × φ ∈ {0.5, 1.5, 3.0}²",                            "9"],
        ["B2  (TISSf)",  "φ(h) = (1/ε₀) exp(−λh),  a = c = 0",              "α × ε₀ × λ  (3 × 2 × 2)",                              "12"],
        ["RL-CBF (ours)","Per-step learned (α, φ, a, c)",                    "—",                                                    "1"],
    ],
    Inches(0.5), Inches(1.95), Inches(12.3), Inches(2.7),
    font_size=14, first_col_bold=True,
    accent_rows={4: GREEN},
)
add_text_block(
    slide,
    "Metric:  combined fall + stuck rate",
    Inches(0.5), Inches(4.9), Inches(12.0), Inches(0.5),
    size=18, bold=True, color=NAVY,
)
add_text_block(
    slide,
    "Fall = body-contact-with-floor termination.  Stuck = horizontal speed below 0.10 m/s for more than half the episode.\n"
    "Including stuck is critical:  earlier evals hid a 22-pp block of failures inside \"timeout successes\" — a robot frozen safely is not a useful robot.",
    Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8),
    size=14, italic=True, color=GRAY,
)

# ---------- Slide 10: In-distribution headline ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Headline result (in-distribution)")
add_text_block(
    slide,
    "RL-CBF achieves the lowest fall, lowest stuck, and lowest combined failure of any of the 25 configurations evaluated.",
    Inches(0.5), Inches(1.15), Inches(12.0), Inches(0.5),
    size=15, italic=True, color=GRAY,
)
add_table(
    slide,
    ["Mode", "Best config", "Fall", "Stuck", "Combined", "h̄", "φ̄"],
    [
        ["B0 plain",       "α = 0.50",                                  "0.293", "0.226", "0.519", "0.305", "0.005"],
        ["B1 ISSf",        "α = 0.50,  φ = 1.50",                       "0.296", "0.119", "0.415", "0.371", "1.50"],
        ["B2 TISSf",       "α = 1.50,  ε₀ = 0.50,  λ = 1.0  (best B)",  "0.243", "0.132", "0.375", "0.371", "1.43"],
        ["RL-CBF (ours)",  "teacher policy",                            "0.231", "0.075", "0.306", "0.413", "0.91"],
    ],
    Inches(0.5), Inches(1.85), Inches(12.3), Inches(2.6),
    font_size=14, first_col_bold=True,
    accent_rows={4: GREEN},
)
add_callout_box(
    slide,
    "+6.9 pp WIN  on combined fall + stuck  vs. the best of 24 hand-tuned baselines",
    Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.9),
    fill_color=GREEN, size=22, bold=True,
)
add_text_block(
    slide,
    "Statistical caveat:  ~134 episodes per config → SE ≈ 3–4 pp; the gap is ~2σ.  Camera-ready run will use paired evaluation (variance ↓ 5–10×).",
    Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.0),
    size=13, italic=True, color=GRAY,
)

# ---------- Slide 11: Behavioral reading ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Why does it win?  Behavioral reading")
add_text_block(
    slide,
    "Hand-tuned baselines that match RL-CBF's fall rate buy that safety with a bigger nominal margin (high φ̄).",
    Inches(0.5), Inches(1.3), Inches(12.0), Inches(0.6),
    size=18, color=BLACK,
)
add_callout_box(
    slide,
    "RL-CBF: same fall safety  +  LOWER φ̄ (0.91)  +  HIGHER h̄ (0.413)",
    Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.0),
    fill_color=NAVY, size=22, bold=True,
)
add_text_block(
    slide,
    "It lets the robot operate closer to its (relaxed) safety boundary while keeping larger absolute distance from obstacles.",
    Inches(0.5), Inches(3.4), Inches(12.0), Inches(0.6),
    size=18, italic=True, color=GRAY,
)
add_text_block(
    slide,
    "The mechanism is adaptation:",
    Inches(0.5), Inches(4.3), Inches(12.0), Inches(0.5),
    size=20, bold=True, color=NAVY,
)
mech = [
    "Spend robustness budget where the scene calls for it.",
    "Save it where it does not — open space, slow-moving obstacles, robot on solid ground.",
    "Static baselines can't do this — they apply the same robustness everywhere.",
]
body = slide.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(12.0), Inches(2.5))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(mech):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = "→  " + b
    r.font.size = Pt(18)
    r.font.color.rgb = BLACK

# ---------- Slide 12: OOD generalization ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Out-of-distribution generalization (5 single-knob shifts)")
add_text_block(
    slide,
    "Each shift pushes ONE randomization axis past the training range while keeping the others at training distribution.",
    Inches(0.5), Inches(1.15), Inches(12.0), Inches(0.5),
    size=14, italic=True, color=GRAY,
)
add_table(
    slide,
    ["Eval axis", "Type", "RL-CBF combined", "Best baseline", "Margin"],
    [
        ["In-dist (ref)",      "mixed",                                 "0.306", "0.375", "+6.9 pp"],
        ["Slippery",           "priv-obs (continuous friction)",        "0.411", "0.453", "+4.1 pp"],
        ["HighDisturbance",    "priv-obs (episodic force / torque)",    "0.341", "0.431", "+9.0 pp"],
        ["HeavyCOM",           "priv-obs (startup COM bias)",           "0.331", "0.390", "+5.9 pp"],
        ["FastObstacles",      "priv-obs (grid history)",               "0.338", "0.444", "+10.5 pp"],
        ["DensePack",          "scene-only (h(x) symmetric)",           "0.338", "0.344", "+0.5 pp  (tie)"],
        ["RealisticCompound",  "all 5 simultaneously",                  "0.500", "0.497", "−0.3 pp  (tie)"],
    ],
    Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.4),
    font_size=13, first_col_bold=True,
    accent_rows={2: LIGHT_GRAY, 4: LIGHT_GRAY, 6: LIGHT_GRAY},
)
add_callout_box(
    slide,
    "5 priv-obs WINS (+4.1 to +10.5 pp)   ·   1 scene-only TIE   ·   1 compositional TIE",
    Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8),
    fill_color=GREEN, size=18, bold=True,
)

# ---------- Slide 13: The pattern ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "The pattern is sharper than \"wins everywhere\"")
add_callout_box(
    slide,
    "RL-CBF wins on every axis whose disturbance signal is in the privileged-observation vector,\n"
    "and ties the best baseline within statistical noise on scene-only shifts.",
    Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.7),
    fill_color=NAVY, size=18, bold=True,
)

add_text_block(
    slide,
    "Wins  (priv-obs aligned)",
    Inches(0.5), Inches(3.3), Inches(6.0), Inches(0.5),
    size=20, bold=True, color=GREEN,
)
wins = [
    "Slippery — friction in dyn block",
    "HighDisturbance — applied force in dyn block",
    "HeavyCOM — COM offset in dyn block",
    "FastObstacles — motion via grid history",
    "(In-dist headline: all of the above)",
]
body = slide.shapes.add_textbox(Inches(0.5), Inches(3.85), Inches(6.0), Inches(3.0))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(wins):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "✓  " + b
    r.font.size = Pt(15)
    r.font.color.rgb = BLACK

add_text_block(
    slide,
    "Ties  (playing field is symmetric)",
    Inches(7.0), Inches(3.3), Inches(6.0), Inches(0.5),
    size=20, bold=True, color=GRAY,
)
ties = [
    "DensePack — only h(x) shifts; baselines see the same h",
    "RealisticCompound — joint high-tail not in training distribution",
]
body = slide.shapes.add_textbox(Inches(7.0), Inches(3.85), Inches(6.0), Inches(3.0))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(ties):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "≈  " + b
    r.font.size = Pt(15)
    r.font.color.rgb = BLACK

add_text_block(
    slide,
    "→ The win comes specifically from privileged information — exactly as the architecture was designed to.",
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7),
    size=16, italic=True, color=NAVY, bold=True,
)

# ---------- Slide 14: Limitations + future work ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Limitations + future work")
add_text_block(
    slide,
    "Sim-only — sim-to-real transfer of a learned safety filter requires three things we have not done",
    Inches(0.5), Inches(1.2), Inches(12.0), Inches(0.5),
    size=18, bold=True, color=NAVY,
)
limits = [
    ("(i)", "Calibrate the LiDAR-derived occupancy grid against real Mid-360 data — drop-out and spurious-return rates we did not inject during training."),
    ("(ii)", "Replace the analytical SDF with a distance-transformed occupancy grid at deploy time (close the train/deploy h-pipeline gap)."),
    ("(iii)", "Student distillation step:  (LiDAR + base_vel + history) → Ẑ ∈ ℝ¹² → frozen π_teacher.  RMA student-side; standard recipe."),
]
body = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(12.3), Inches(3.5))
tf = body.text_frame
tf.word_wrap = True
for i, (lead, b) in enumerate(limits):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r1 = p.add_run()
    r1.text = lead + "  "
    r1.font.size = Pt(18)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = b
    r2.font.size = Pt(16)
    r2.font.color.rgb = BLACK

add_text_block(
    slide,
    "Other open items",
    Inches(0.5), Inches(5.3), Inches(12.0), Inches(0.5),
    size=18, bold=True, color=NAVY,
)
others = [
    "b slot reserved but unused  →  would require SOCP solver instead of closed-form half-space projection.",
    "Compositional generalization (RealisticCompound tie)  →  joint high-tail untrained.",
    "Paired evaluation (same seed list across all 25 configs)  →  variance reduction 5–10× for camera-ready.",
]
body = slide.shapes.add_textbox(Inches(0.7), Inches(5.85), Inches(12.0), Inches(2.0))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(others):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = "•  " + b
    r.font.size = Pt(14)
    r.font.color.rgb = GRAY

# ---------- Slide 15: Summary / Q&A ----------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "Summary")
add_callout_box(
    slide,
    "A privileged-observation RL teacher learns the four robust-CBF parameters per timestep.\n"
    "Beats the best of 24 hand-tuned baselines by +6.9 pp combined fall + stuck in-distribution.\n"
    "Wins on every privileged-obs aligned OOD axis.  Ties on scene-only / compositional shifts.",
    Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.4),
    fill_color=NAVY, size=18, bold=False,
)
add_text_block(
    slide,
    "Why this matters",
    Inches(0.5), Inches(4.0), Inches(12.0), Inches(0.5),
    size=20, bold=True, color=NAVY,
)
why = [
    "Robust-CBF parameters can be adapted online — no need to commit to one conservative tuple.",
    "Drop-in design:  works with any frozen locomotion stack.  Safety guarantee preserved by output non-negativity.",
    "Per-step adaptation gives a clean differentiation from per-episode α-calibration baselines (e.g., SHIELD).",
]
body = slide.shapes.add_textbox(Inches(0.5), Inches(4.55), Inches(12.3), Inches(2.5))
tf = body.text_frame
tf.word_wrap = True
for i, b in enumerate(why):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run()
    r.text = "→  " + b
    r.font.size = Pt(16)
    r.font.color.rgb = BLACK

add_text_block(
    slide,
    "Questions?",
    Inches(0.5), Inches(6.7), Inches(12.0), Inches(0.6),
    size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
)

# ---------- Save ----------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")
